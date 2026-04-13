from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
RED = RGBColor(0xDC, 0x14, 0x3C)       # Lenovo-ish red
DARK_RED = RGBColor(0xB0, 0x10, 0x30)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MID_GRAY = RGBColor(0x66, 0x66, 0x66)
ACCENT_RED_BG = RGBColor(0xFD, 0xE8, 0xEB)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_red_header_bar(slide, title_text):
    """Add a red bar across the top with white title text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(1.1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RED
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.8)
    tf.margin_top = Inches(0.15)


def add_bottom_bar(slide):
    """Add a thin red accent bar at the bottom."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.1), SLIDE_W, Inches(0.4)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RED
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Bsharp Reco  |  AI-Powered Retail Recommendations"
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    tf.margin_top = Inches(0.05)


def add_white_bg(slide):
    """Set slide background to white."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_screenshot_placeholder(slide, left, top, width, height):
    """Add a dashed-border placeholder box for a screenshot."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_GRAY
    shape.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    shape.line.width = Pt(1.5)
    shape.line.dash_style = 4  # dash

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    p = tf.paragraphs[0]
    p.text = "[ Screenshot Placeholder ]"
    p.font.size = Pt(14)
    p.font.color.rgb = MID_GRAY
    p.font.italic = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.word_wrap = True
    from pptx.enum.text import MSO_ANCHOR
    shape.text_frame.auto_size = None
    # Vertically center
    shape.text_frame.paragraphs[0].space_before = Pt(round(height / Pt(1) * 0.35))


def add_bullet_textbox(slide, left, top, width, height, bullets, font_size=16, bold_prefix=True):
    """Add a text box with bullet points. Bullets can be strings or (bold_part, rest) tuples."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.space_after = Pt(8)
        p.space_before = Pt(4)
        p.level = 0

        if isinstance(bullet, tuple):
            # (bold_text, normal_text)
            run_bold = p.add_run()
            run_bold.text = bullet[0]
            run_bold.font.size = Pt(font_size)
            run_bold.font.bold = True
            run_bold.font.color.rgb = DARK_GRAY

            run_normal = p.add_run()
            run_normal.text = bullet[1]
            run_normal.font.size = Pt(font_size)
            run_normal.font.bold = False
            run_normal.font.color.rgb = MID_GRAY
        else:
            run = p.add_run()
            run.text = bullet
            run.font.size = Pt(font_size)
            run.font.color.rgb = DARK_GRAY

    return txBox


def add_red_number_badge(slide, left, top, number):
    """Add a red circle with a number in it."""
    size = Inches(0.45)
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RED
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    tf.margin_top = Inches(0.02)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)


# ============================================================
# SLIDE 1: Title Slide
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_white_bg(slide1)

# Big red block in center
shape = slide1.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(4.2)
)
shape.fill.solid()
shape.fill.fore_color.rgb = RED
shape.line.fill.background()

# Title
txBox = slide1.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11), Inches(1.2))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Bsharp Reco"
p.font.size = Pt(52)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.LEFT

# Subtitle
txBox2 = slide1.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11), Inches(1.2))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.text = "AI-Powered Retail Recommendation Platform"
p2.font.size = Pt(28)
p2.font.color.rgb = WHITE
p2.alignment = PP_ALIGN.LEFT

# Bottom section with key tagline
txBox4 = slide1.shapes.add_textbox(Inches(0.8), Inches(4.8), Inches(11), Inches(1.5))
tf4 = txBox4.text_frame
tf4.word_wrap = True
p4 = tf4.paragraphs[0]
p4.text = "From lack of clarity to confident choice —"
p4.font.size = Pt(22)
p4.font.color.rgb = DARK_GRAY
p4.alignment = PP_ALIGN.LEFT

p5 = tf4.add_paragraph()
p5.text = "turning in-store interactions into qualified retail leads"
p5.font.size = Pt(22)
p5.font.color.rgb = DARK_GRAY
p5.alignment = PP_ALIGN.LEFT

add_bottom_bar(slide1)


# ============================================================
# SLIDE 2: System Objective
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide2)
add_red_header_bar(slide2, "System Objective")
add_bottom_bar(slide2)

# Content - left side
bullets_obj = [
    ("Guide customers from uncertainty to confident choice", " — through AI-driven consultative selling at the point of sale"),
    ("Combine voice discovery, adaptive questioning & explainable recommendations", " into a single seamless guided experience"),
    ("Capture every interaction as a qualified retail lead", " — with full preference context, journey data, and product interest for CRM follow-up"),
    ("Built for scale", " — deploy once, run across 8,000+ retail outlets with centralized control over products, recommendations, and business rules"),
]
add_bullet_textbox(slide2, Inches(0.8), Inches(1.5), Inches(6.5), Inches(5.0), bullets_obj, font_size=16)

# Screenshot on right — ss1 (Discovery mode) — width only, preserve aspect ratio
slide2.shapes.add_picture("ss1.png", Inches(7.2), Inches(1.8), width=Inches(5.8))


# ============================================================
# SLIDE 3: Top 5 Features
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide3)
add_red_header_bar(slide3, "Top 5 Features")
add_bottom_bar(slide3)

features = [
    ("AI Voice & Text Discovery", "Customers speak or type their needs; AI transcribes and extracts preference tags instantly. Supports 90+ languages."),
    ("LLM-Powered Adaptive Guided Questions", "No pre-authored question bank. AI generates personalized questions in real-time based on customer input, product catalog, and session context until a confident recommendation is served."),
    ("Intelligent Recommendations with Explainability", "Weighted scoring engine produces top 3 product picks, each with a clear \"Why this product?\" explanation grounded in matched features and benefits."),
    ("Brand Moderation & Business Rules Engine", "Upload plain-English moderation rules; business overrides enable campaign boosts, SKU pushes, and product suppressions — all without code changes."),
    ("Automated Lead Capture & Email Handoff", "Customer preferences, selected product, and full journey context captured as a qualified CRM lead. One-click email share to customer."),
]

y_pos = Inches(1.4)
for i, (title, desc) in enumerate(features):
    # Number badge
    add_red_number_badge(slide3, Inches(0.7), y_pos + Inches(0.05), i + 1)

    # Feature card
    card = slide3.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.3), y_pos, Inches(11.2), Inches(0.95)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = ACCENT_RED_BG if i % 2 == 0 else LIGHT_GRAY
    card.line.fill.background()

    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)

    p_title = tf.paragraphs[0]
    run_t = p_title.add_run()
    run_t.text = title
    run_t.font.size = Pt(15)
    run_t.font.bold = True
    run_t.font.color.rgb = DARK_RED

    p_desc = tf.add_paragraph()
    run_d = p_desc.add_run()
    run_d.text = desc
    run_d.font.size = Pt(12)
    run_d.font.color.rgb = MID_GRAY

    y_pos += Inches(1.08)


# ============================================================
# SLIDE 4: Benefits of These Features
# ============================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide4)
add_red_header_bar(slide4, "Benefits for Lenovo")
add_bottom_bar(slide4)

benefits = [
    ("Voice & Text Discovery",
     "Reduces time-to-engagement; removes friction for customers who can't articulate specs"),
    ("Adaptive Guided Questions",
     "Replicates an expert salesperson's consultative approach at every outlet, consistently"),
    ("Explainable Recommendations",
     "Builds customer trust; gives store staff ready-made talking points for closing"),
    ("Moderation & Business Rules",
     "Control what gets recommended — push new launches, suppress EOL products, enforce regional rules in real-time"),
    ("Lead Capture & Email Handoff",
     "Every walk-in becomes a trackable, follow-up-ready lead with complete preference data for your CRM pipeline"),
]

# Table-like layout
y_start = Inches(1.5)
for i, (feature, benefit) in enumerate(benefits):
    bg_color = ACCENT_RED_BG if i % 2 == 0 else LIGHT_GRAY

    # Feature label (left column)
    left_box = slide4.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y_start, Inches(3.5), Inches(0.95)
    )
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = RED
    left_box.line.fill.background()
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    tf_left.margin_left = Inches(0.2)
    tf_left.margin_top = Inches(0.15)
    p_l = tf_left.paragraphs[0]
    run_l = p_l.add_run()
    run_l.text = feature
    run_l.font.size = Pt(15)
    run_l.font.bold = True
    run_l.font.color.rgb = WHITE
    p_l.alignment = PP_ALIGN.CENTER

    # Benefit text (right column)
    right_box = slide4.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.3), y_start, Inches(8.3), Inches(0.95)
    )
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = bg_color
    right_box.line.fill.background()
    tf_right = right_box.text_frame
    tf_right.word_wrap = True
    tf_right.margin_left = Inches(0.3)
    tf_right.margin_top = Inches(0.15)
    p_r = tf_right.paragraphs[0]
    run_r = p_r.add_run()
    run_r.text = benefit
    run_r.font.size = Pt(14)
    run_r.font.color.rgb = DARK_GRAY

    y_start += Inches(1.08)


# ============================================================
# SLIDE 5: Vision — From Browsing to Guided Selling
# ============================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide5)
add_red_header_bar(slide5, "Vision: From Browsing to Guided Selling")
add_bottom_bar(slide5)

# Left content
vision_bullets = [
    ("Traditional retail: ", "Customer browses → overwhelmed by options → leaves or picks randomly"),
    ("Bsharp Reco: ", "Customer speaks → AI understands → guided conversation → confident choice"),
    ("Paradigm shift: ", "From \"search & browse\" to \"consult & recommend\""),
    ("Data-driven retail: ", "Every interaction generates actionable insight for Lenovo's retail strategy"),
]
add_bullet_textbox(slide5, Inches(0.8), Inches(1.5), Inches(6.2), Inches(5.0), vision_bullets, font_size=16)

# Screenshot on right — ss2 (Recommendations) — width only, preserve aspect ratio
slide5.shapes.add_picture("ss2.png", Inches(7.0), Inches(1.8), width=Inches(5.8))


# ============================================================
# SLIDE 6: Capability — Self-Learning AI Engine
# ============================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide6)
add_red_header_bar(slide6, "Capability: Self-Learning AI Engine")
add_bottom_bar(slide6)

ai_bullets = [
    ("Side-by-Side Product Comparison", " — customers can compare any two recommended products with AI-generated commentary highlighting key differences, trade-offs, and best-fit scenarios."),
    ("RAG-Grounded Product Chatbot", " — answers real-time customer questions from actual specs and documentation, not hallucinations."),
    ("Feedback-Driven Learning Loop", " — star ratings analyzed nightly; system learns which recommendation patterns work and avoids ones that don't."),
    ("Real-Time Recommendation Control", " — push new launches, suppress end-of-life products, and run campaigns instantly without downtime."),
]
add_bullet_textbox(slide6, Inches(0.8), Inches(1.5), Inches(6.2), Inches(5.0), ai_bullets, font_size=16)

# Screenshot on right — ss3 (Comparison) — width only, preserve aspect ratio
slide6.shapes.add_picture("ss3.png", Inches(7.0), Inches(1.8), width=Inches(5.8))


# ============================================================
# SLIDE 7: Capability — Enterprise-Ready at Scale
# ============================================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide7)
add_red_header_bar(slide7, "Capability: Enterprise-Ready at Scale")
add_bottom_bar(slide7)

enterprise_bullets = [
    ("Cost-efficient deployment", " — runs on CPU-only infrastructure, no expensive GPU needed"),
    ("Seamless integration", " — plugs into your existing systems and user management with zero disruption"),
    ("No-code admin control", " — manage products, moderation rules, scoring configs, and campaigns through an intuitive interface"),
    ("Rich analytics from every interaction", " — customer preferences, drop-off points, top-performing products, and conversion patterns across all outlets"),
    ("Full audit trail", " — every session logged for compliance, dispute resolution, and performance analysis"),
]
add_bullet_textbox(slide7, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.0), enterprise_bullets, font_size=16)


# ============================================================
# SAVE
# ============================================================
output_path = "/Users/tanujsadasivam/Desktop/Bsharp_Reco_Implementation/demoppt/Bsharp_Reco_Lenovo.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
