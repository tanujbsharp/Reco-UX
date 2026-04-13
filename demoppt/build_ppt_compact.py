from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Colors
RED = RGBColor(0xDC, 0x14, 0x3C)
DARK_RED = RGBColor(0xB0, 0x10, 0x30)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MID_GRAY = RGBColor(0x66, 0x66, 0x66)
ACCENT_RED_BG = RGBColor(0xFD, 0xE8, 0xEB)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = Inches(13.333)


def add_red_header_bar(slide, title_text):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(1.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RED
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.8)
    tf.margin_top = Inches(0.12)


def add_bottom_bar(slide):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.1), SLIDE_W, Inches(0.4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RED
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Bsharp Reco  |  AI-Powered Retail Recommendations"
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    tf.margin_top = Inches(0.05)


def add_white_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_red_number_badge(slide, left, top, number):
    size = Inches(0.38)
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RED
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    tf.margin_top = Inches(0.01)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)


# ============================================================
# SLIDE 1: Title — ss1 background with text overlay
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide1)

# ss1 image centered at top
slide1.shapes.add_picture("ss1.png", Inches(2.9), Inches(0.3), width=Inches(7.5))

# Title below image
txBox = slide1.shapes.add_textbox(Inches(0), Inches(4.2), SLIDE_W, Inches(1.0))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Bsharp Reco"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = RED
p.alignment = PP_ALIGN.CENTER

# Tagline below title
txBox2 = slide1.shapes.add_textbox(Inches(1.5), Inches(5.2), Inches(10.3), Inches(1.5))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.text = "From lack of clarity to confident choice —"
p2.font.size = Pt(22)
p2.font.color.rgb = DARK_GRAY
p2.alignment = PP_ALIGN.CENTER
p3 = tf2.add_paragraph()
p3.text = "turning in-store interactions into qualified retail leads"
p3.font.size = Pt(22)
p3.font.color.rgb = MID_GRAY
p3.alignment = PP_ALIGN.CENTER

add_bottom_bar(slide1)


# ============================================================
# SLIDE 2: Features (3 items, concise)
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide2)
add_red_header_bar(slide2, "What It Does")
add_bottom_bar(slide2)

features = [
    ("AI-Powered Discovery & Guided Questions",
     "Voice or text input — AI extracts preferences, generates adaptive questions, and narrows down the right product."),
    ("Explainable Recommendations",
     "Top 3 picks with clear \"Why this product?\" reasoning grounded in matched features and benefits."),
    ("Brand Moderation & Business Rules",
     "Plain-English rules to push launches, suppress products, and run campaigns — no code changes needed."),
]

y_pos = Inches(1.5)
for i, (title, desc) in enumerate(features):
    add_red_number_badge(slide2, Inches(0.6), y_pos + Inches(0.12), i + 1)

    card = slide2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.1), y_pos, Inches(6.0), Inches(1.2)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = ACCENT_RED_BG if i % 2 == 0 else LIGHT_GRAY
    card.line.fill.background()

    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0.08)

    p_title = tf.paragraphs[0]
    run_t = p_title.add_run()
    run_t.text = title
    run_t.font.size = Pt(15)
    run_t.font.bold = True
    run_t.font.color.rgb = DARK_RED

    p_desc = tf.add_paragraph()
    run_d = p_desc.add_run()
    run_d.text = desc
    run_d.font.size = Pt(12)
    run_d.font.color.rgb = MID_GRAY

    y_pos += Inches(1.45)

# Screenshot — ss2 (recommendations)
slide2.shapes.add_picture("ss2.png", Inches(7.3), Inches(1.3), width=Inches(5.6))


# ============================================================
# SLIDE 3: Capabilities
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide3)
add_red_header_bar(slide3, "Platform Capabilities")
add_bottom_bar(slide3)

capabilities = [
    ("Product Details & Side-by-Side Comparison",
     "View full specs of any product, or compare two products with AI-generated commentary on differences and best-fit scenarios."),
    ("Self-Learning Recommendations",
     "Star ratings analyzed nightly — system learns what works and stops recommending what doesn't."),
    ("Real-Time Campaign Control",
     "Push new launches, suppress EOL products, adjust scoring — instantly, no downtime."),
    ("Analytics Across All Outlets",
     "Customer preferences, conversion patterns, and top-performing products — all in one view."),
    ("Lead Capture & Handoff",
     "Every interaction captured as a qualified CRM lead with full journey context and one-click email share."),
]

y_pos = Inches(1.25)
for i, (title, desc) in enumerate(capabilities):
    add_red_number_badge(slide3, Inches(0.5), y_pos + Inches(0.08), i + 1)

    card = slide3.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), y_pos, Inches(5.8), Inches(0.95)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = ACCENT_RED_BG if i % 2 == 0 else LIGHT_GRAY
    card.line.fill.background()

    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.05)

    p_title = tf.paragraphs[0]
    run_t = p_title.add_run()
    run_t.text = title
    run_t.font.size = Pt(14)
    run_t.font.bold = True
    run_t.font.color.rgb = DARK_RED

    p_desc = tf.add_paragraph()
    run_d = p_desc.add_run()
    run_d.text = desc
    run_d.font.size = Pt(11)
    run_d.font.color.rgb = MID_GRAY

    y_pos += Inches(1.05)

# Screenshot — ss3 (comparison)
slide3.shapes.add_picture("ss3.png", Inches(7.2), Inches(1.3), width=Inches(5.7))


# ============================================================
# SAVE
# ============================================================
output_path = "/Users/tanujsadasivam/Desktop/Bsharp_Reco_Implementation/demoppt/Bsharp_Reco_Lenovo_Compact.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
